"""Generate an editable PowerPoint of the GAPGen system overview.

Minimal style, no title. Boxes = modules, arrows = data flow. The three
novelties (gameplay-ontology planning, multi-reference conditioning, hybrid
gate + human-in-the-loop) are highlighted in accent color with a small tag.

Run: python paper/figures/make_system_overview.py
Output: paper/figures/system_overview.pptx  (fully editable in PowerPoint/Keynote)
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT = pathlib.Path(__file__).parent / 'system_overview.pptx'

# Minimal monochrome + one accent for novelty.
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x88, 0x88, 0x88)
LINE = RGBColor(0x55, 0x55, 0x55)
FILL = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE8, 0x55, 0x00)       # novelty stroke / tag
ACCENT_FILL = RGBColor(0xFF, 0xF3, 0xEA)  # novelty light fill
IO_FILL = RGBColor(0xF2, 0xF2, 0xF2)      # input/output boxes

FONT = 'Arial'


def _no_autofit(tf):
    # keep text size fixed; user edits freely
    tf.word_wrap = True


def box(slide, x, y, w, h, lines, *, novelty=False, io=False, fs=13):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT_FILL if novelty else (IO_FILL if io else FILL)
    shp.line.color.rgb = ACCENT if novelty else LINE
    shp.line.width = Pt(1.75 if novelty else 1.0)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _no_autofit(tf)
    if isinstance(lines, str):
        lines = [(lines, True)]
    for i, item in enumerate(lines):
        text, bold = item if isinstance(item, tuple) else (item, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(fs if bold else fs - 2)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = INK if bold else MUTED
    return shp


def tag(slide, x, y, text):
    """Small accent 'novelty' pill."""
    w, h = 1.15, 0.26
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame; _no_autofit(tf)
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.bold = True
    r.font.name = FONT; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shp


def _arrow_head(conn, color, dashed=False):
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    if dashed:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)


def arrow(slide, x1, y1, x2, y2, *, color=LINE, dashed=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.shadow.inherit = False
    _arrow_head(conn, color, dashed)
    return conn


def label(slide, x, y, w, text, *, color=MUTED, fs=10, align=PP_ALIGN.CENTER, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame; _no_autofit(tf)
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.name = FONT; r.font.color.rgb = color; r.font.italic = italic
    return tb


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ---- Row geometry ----
    mid_y = 3.05          # vertical center band for the main pipeline row
    bh = 1.5              # main box height
    top = mid_y

    # Column x/width
    x_in = 0.45; w_in = 1.85
    x_plan = 2.75; w_plan = 3.15
    x_gen = 6.35; w_gen = 2.35
    x_gate = 9.15; w_gate = 2.55
    x_out = 12.05; w_out = 1.05  # narrow; will widen output as tall box

    # 1) INPUT
    box(slide, x_in, top, w_in, bh, [
        ('Inputs', True),
        ('Original sprites', False),
        ('Gameplay-role ontology', False),
        ('(asset_roles.json)', False),
    ], io=True)

    # 2) ONTOLOGY-DRIVEN PLANNERS  (novelty 1)
    tag(slide, x_plan + w_plan/2 - 0.575, top - 0.34, 'NOVELTY 1')
    box(slide, x_plan, top, w_plan, bh, [
        ('Ontology-driven planners', True),
        ('Style · Theme · Family-style · Stage', False),
        ('roles \u2192 concrete visual specs', False),
    ], novelty=True)

    # 3) GENERATOR
    box(slide, x_gen, top, w_gen, bh, [
        ('Image generator', True),
        ('(diffusion VLM)', False),
        ('prompt + reference set', False),
    ])

    # 4) HYBRID GATES + ITERATE  (novelty 3)
    tag(slide, x_gate + w_gate/2 - 0.575, top - 0.34, 'NOVELTY 3')
    box(slide, x_gate, top, w_gate, bh, [
        ('Hybrid validation', True),
        ('Postprocess (rule)  +  VLM critic', False),
        ('pass / needs-review / retry', False),
    ], novelty=True)

    # 5) OUTPUT (tall, right)
    box(slide, x_out, top, w_out + 0.2, bh, [
        ('Staged pack', True),
        ('pass', False),
        ('needs-', False),
        ('review', False),
    ], io=True, fs=12)

    # ---- Main horizontal data flow arrows ----
    cy = top + bh/2
    arrow(slide, x_in + w_in, cy, x_plan, cy)
    arrow(slide, x_plan + w_plan, cy, x_gen, cy)
    arrow(slide, x_gen + w_gen, cy, x_gate, cy)
    arrow(slide, x_gate + w_gate, cy, x_out, cy)

    # ---- Iterate-with-feedback loop (gate -> generator), accent dashed ----
    loop_y = top + bh + 0.55
    # down from gate
    arrow(slide, x_gate + 0.5, top + bh, x_gate + 0.5, loop_y, color=ACCENT)
    # left back under generator
    arrow(slide, x_gate + 0.5, loop_y, x_gen + w_gen/2, loop_y, color=ACCENT)
    # up into generator
    arrow(slide, x_gen + w_gen/2, loop_y, x_gen + w_gen/2, top + bh, color=ACCENT)
    label(slide, x_gen + 0.2, loop_y + 0.02, x_gate - x_gen + 1.0,
          'iterate with fix instructions  (\u2264 N)', color=ACCENT, fs=10,
          align=PP_ALIGN.CENTER)

    # ---- Reference set feeding the generator  (novelty 2) ----
    ref_y = top - 1.75
    ref_h = 1.25
    refx = x_gen - 0.55; refw = w_gen + 1.1
    tag(slide, refx + refw/2 - 0.575, ref_y - 0.34, 'NOVELTY 2')
    box(slide, refx, ref_y, refw, ref_h, [
        ('Multi-reference conditioning', True),
        ('A: original \u2192 function', False),
        ('family anchor \u2192 style cohesion', False),
        ('prev stage \u2192 progression chain', False),
    ], novelty=True, fs=12)
    # arrow down into generator
    arrow(slide, x_gen + w_gen/2, ref_y + ref_h, x_gen + w_gen/2, top, color=ACCENT)
    # planners also inform the reference/anchor plan (light link)
    arrow(slide, x_plan + w_plan, top + 0.15, refx, ref_y + ref_h - 0.2,
          color=MUTED, dashed=True)

    # ---- Human-in-the-loop (final say) below output ----
    human_y = top + bh + 0.9
    hb = box(slide, x_gate + 0.1, human_y, w_gate + 1.0, 0.85, [
        ('Human review \u2192 final say', True),
        ('accept / edit / regenerate', False),
    ], io=True, fs=12)
    arrow(slide, x_out + (w_out + 0.2)/2, top + bh, x_out + (w_out + 0.2)/2, human_y,
          color=LINE)

    # ---- One-line caption of the contribution (bottom) ----
    label(slide, 0.5, 6.85, 12.3,
          'Gameplay-consistent asset generation: ontology-planned specs + multi-reference '
          'conditioning + hybrid rule/VLM gates keep function, cohesion and progression under human control.',
          color=MUTED, fs=11, align=PP_ALIGN.CENTER, italic=True)

    prs.save(OUT)
    return OUT


if __name__ == '__main__':
    p = build()
    print(f'wrote {p}')
